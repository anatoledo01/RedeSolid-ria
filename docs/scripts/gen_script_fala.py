"""
Gera o script de fala da apresentação.

Saídas:
  1. Apresentacao_RedeSolidaria.pptx — atualizado com speaker notes
     (aparecem no Modo Apresentador do PowerPoint/Keynote)
  2. Script_Apresentacao.pdf — versão imprimível agrupada por
     apresentador, com tempo estimado, dica visual e fala completa.
"""

import os
from pptx import Presentation
from pptx.util import Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT, TA_CENTER
from reportlab.lib import colors
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                 PageBreak, Table, TableStyle)


ROOT = "/Users/jhonatanrabelo/Documents/Fatec/Eng 3/Projeto/RedeSolid-ria"
PPTX_PATH = os.path.join(ROOT, "Apresentacao_RedeSolidaria.pptx")
PDF_PATH  = os.path.join(ROOT, "Script_Apresentacao.pdf")

# Cores p/ o PDF (mesmas do PPTX)
NAVY  = colors.HexColor("#1f3a5f")
GREEN = colors.HexColor("#059669")
BLUE  = colors.HexColor("#2563eb")
PINK  = colors.HexColor("#db2777")
GREY  = colors.HexColor("#475569")
GREY_L = colors.HexColor("#f1f5f9")
TEXT  = colors.HexColor("#0f172a")

PRESENTERS = {
    "Ana":      {"name": "Ana Toledo",      "color": PINK},
    "Vitor":    {"name": "Vitor Campos",    "color": BLUE},
    "Jhonatan": {"name": "Jhonatan Rabelo", "color": GREEN},
    "Grupo":    {"name": "Grupo",           "color": NAVY},
}


# =============================================================
# CONTEÚDO — uma entrada por slide, na ordem do .pptx
# Campos: presenter, title, time, cue, speech
# =============================================================
SCRIPT = [
    {
        "presenter": "Ana",
        "title": "Capa — abertura",
        "time": "20s",
        "cue": "Slide de capa, antes de virar para a Agenda.",
        "speech": (
            "Boa tarde a todos. Somos do grupo da disciplina de "
            "Engenharia de Software III, e hoje vamos apresentar a "
            "Rede Solidária, nossa plataforma de gestão de doações "
            "desenvolvida em parceria com a ONG Cidade Social de Mogi. "
            "Eu sou a Ana Toledo, e comigo apresentam o Vitor Campos e "
            "o Jhonatan Rabelo. O Edilson Junior, que faz parte do grupo, "
            "foi quem nos conectou com a ONG."
        ),
    },
    {
        "presenter": "Ana",
        "title": "Agenda e divisão",
        "time": "45s",
        "cue": "Aponte para os 3 blocos coloridos da esquerda.",
        "speech": (
            "Para facilitar, dividimos a apresentação em três blocos. "
            "Eu começo contando o problema que encontramos na ONG e a "
            "solução que propusemos. Em seguida o Vitor entra na "
            "modelagem do sistema, mostrando os diagramas. Por fim, o "
            "Jhonatan apresenta a implementação, os padrões de projeto "
            "que aplicamos, o cronograma e a conclusão. Vocês vão ver "
            "uma chip no canto superior direito de cada slide indicando "
            "quem está apresentando."
        ),
    },
    {
        "presenter": "Ana",
        "title": "O Problema — Parte 1 (Contexto)",
        "time": "1min",
        "cue": "Aponte para a citação destacada no quadro verde.",
        "speech": (
            "Nosso projeto nasceu de uma necessidade real. A Cidade "
            "Social de Mogi é uma ONG que monta e distribui cestas "
            "básicas para famílias em situação de vulnerabilidade aqui "
            "em Mogi das Cruzes. Conversamos com o Jonas Oliveira, que "
            "é o principal stakeholder, e ele descreveu o cenário "
            "atual com essa frase que você vê aqui em destaque: hoje "
            "se perde produto por validade vencida, e não dá pra saber "
            "direito quem doou o quê nem para qual família foi cada "
            "cesta. Tudo é controlado no caderno."
        ),
    },
    {
        "presenter": "Ana",
        "title": "O Problema — Parte 2 (Dores específicas)",
        "time": "1min15s",
        "cue": "Três cards: leia o título de cada um e amplie um pouco.",
        "speech": (
            "A partir dessa conversa, identificamos três dores claras. "
            "Primeiro, controle de estoque: a ONG precisa monitorar a "
            "validade dos produtos pra montar cestas equilibradas com "
            "o que está em dia. Segundo, cadastro de pessoas: precisa "
            "manter registro completo de quem doa e de quem recebe, "
            "com CPF, nome e contato. E terceiro, rastreamento das "
            "cestas: saber quando foi montada, quem recebeu e a "
            "validade. Resumindo: precisamos transformar caderno em "
            "um processo digital, simples e auditável."
        ),
    },
    {
        "presenter": "Ana",
        "title": "A Solução",
        "time": "1min",
        "cue": "Aponte para os 4 cards verdes (papéis).",
        "speech": (
            "Nossa resposta foi a Rede Solidária, uma plataforma web "
            "full-stack que conecta quatro papéis em torno do ciclo "
            "da doação. O Doador publica os itens e acompanha o "
            "destino. O Recebedor busca o que precisa e reserva. O "
            "Voluntário aceita as entregas e atualiza o status. E o "
            "Administrador aprova usuários, modera, e tem acesso ao "
            "dashboard e aos logs de auditoria. O fluxo é simples: o "
            "Doador publica, o Recebedor reserva, o Voluntário "
            "entrega, e os dois lados se avaliam. Agora passo a "
            "palavra para o Vitor, que vai mostrar como modelamos "
            "isso tudo."
        ),
    },
    {
        "presenter": "Vitor",
        "title": "Transição → Vitor",
        "time": "15s",
        "cue": "Slide cheio na cor azul; entre rápido.",
        "speech": (
            "Obrigado, Ana. Eu vou ficar com o bloco de modelagem do "
            "sistema. Vamos começar pela stack que escolhemos, e em "
            "seguida olhar os cinco diagramas que descrevem o sistema."
        ),
    },
    {
        "presenter": "Vitor",
        "title": "Stack Tecnológica",
        "time": "1min",
        "cue": "Leia as quatro camadas (Frontend, Backend, Banco, Dev).",
        "speech": (
            "Para a stack, no frontend usamos Next.js 15 com React 19 "
            "e TypeScript — bem moderno, com Tailwind v4 e React "
            "Query para o estado de dados. No backend, NestJS 11 com "
            "Prisma como ORM e autenticação por JWT com Passport. "
            "O banco é PostgreSQL 16, rodando em Docker. E pra "
            "desenvolvimento usamos Docker Compose e ts-node. O "
            "backend gera Swagger automático, que dá pra acessar em "
            "localhost 3001 barra api barra docs."
        ),
    },
    {
        "presenter": "Vitor",
        "title": "IDEF0 N0 — Contexto",
        "time": "1min15s",
        "cue": "Aponte: ENTRADAS à esquerda, CONTROLES em cima, SAÍDAS "
                "à direita, MECANISMOS embaixo.",
        "speech": (
            "Esse é o IDEF0 de nível zero — a visão de contexto. A "
            "caixa central representa a plataforma inteira. À "
            "esquerda estão as entradas: itens para doação, "
            "necessidades de recebedores e os dados de cadastro. Em "
            "cima ficam os controles: LGPD, as regras de negócio "
            "como os papéis e os status, e a política de aprovação. "
            "À direita saem os resultados: doações entregues, "
            "avaliações, relatórios e logs. E embaixo estão os "
            "mecanismos que tornam isso possível: os próprios "
            "usuários, o frontend Next, o backend Nest e o banco "
            "PostgreSQL."
        ),
    },
    {
        "presenter": "Vitor",
        "title": "IDEF0 N1 — Decomposição",
        "time": "1min15s",
        "cue": "Aponte para cada caixa A1→A5 em sequência.",
        "speech": (
            "Quando decompomos esse contexto, chegamos em cinco "
            "subatividades em escada. A1 cuida de Usuários e "
            "Autenticação — é onde o login produz o token JWT. Esse "
            "token alimenta A2, o catálogo de doações. Quando uma "
            "doação é reservada, dispara A3, a coordenação das "
            "entregas pelo voluntário. Quando a entrega termina, "
            "vem A4 com as avaliações. E A5, no canto, é a "
            "administração e auditoria — consolida tudo em "
            "relatórios e logs. Cada caixa corresponde a um "
            "Module do nosso NestJS."
        ),
    },
    {
        "presenter": "Vitor",
        "title": "Diagrama de Casos de Uso",
        "time": "1min15s",
        "cue": "Bonecos à esquerda (Doador, Recebedor, Voluntário), "
                "Admin à direita.",
        "speech": (
            "Aqui temos os casos de uso. Os três atores comuns "
            "estão à esquerda — Doador, Recebedor e Voluntário — e "
            "o Administrador, que tem permissões diferentes, está "
            "isolado à direita, em laranja. Todos passam por "
            "cadastrar-se e autenticar. Daí cada um tem suas "
            "operações: o Doador publica e cancela, o Recebedor "
            "busca e reserva, e o Voluntário aceita entregas e "
            "atualiza status. Repare também nas relações com "
            "estereótipo: publicar doação inclui anexar fotos, e "
            "aceitar entrega estende reservar."
        ),
    },
    {
        "presenter": "Vitor",
        "title": "DFD Nível 1",
        "time": "1min15s",
        "cue": "Externas à esquerda/direita, processos P1–P6 no centro, "
                "depósitos D1–D5 embaixo.",
        "speech": (
            "Esse é o DFD de nível um, na notação Gane e Sarson. "
            "As entidades externas amarelas são os atores. No "
            "centro, em verde, estão os processos do sistema, "
            "numerados de P1 a P6. E embaixo, em azul, os "
            "depósitos de dados, que são as nossas tabelas — "
            "Usuários, Doações, Entregas, Reviews e o AuditLog. "
            "As setas mostram que fluxo de informação atravessa "
            "cada processo. Por exemplo: o doador entrega fotos e "
            "descrição para P2; P2 grava em D2; quando reserva, "
            "a doação passa para P4 que cria a entrega em D3."
        ),
    },
    {
        "presenter": "Vitor",
        "title": "DER — Modelo de Dados",
        "time": "1min15s",
        "cue": "User no topo, Donation/DeliveryRequest/Review no meio, "
                "leaves embaixo.",
        "speech": (
            "E aqui o DER, que reflete exatamente o nosso schema "
            "no Prisma. Em cima a entidade User, que é central — "
            "todos os papéis vivem na mesma tabela e se "
            "diferenciam pelo campo role. Na linha do meio, as três "
            "entidades principais do ciclo: Donation, "
            "DeliveryRequest e Review. E embaixo as auxiliares: "
            "Category, DonationImage, Address, RefreshToken e "
            "AuditLog. As cardinalidades em vermelho mostram, por "
            "exemplo, que um Usuário pode ter várias Doações, "
            "mas uma Doação tem só um Doador. Agora passo pro "
            "Jhonatan, que vai falar de como isso virou código."
        ),
    },
    {
        "presenter": "Jhonatan",
        "title": "Transição → Jhonatan",
        "time": "15s",
        "cue": "Slide cheio na cor verde; entre rápido.",
        "speech": (
            "Valeu, Vitor. Eu fico com o bloco de implementação. "
            "Vou mostrar como traduzimos esses diagramas em "
            "código orientado a objetos, os três padrões de "
            "projeto que aplicamos no backend, e fechar com o "
            "cronograma e a conclusão."
        ),
    },
    {
        "presenter": "Jhonatan",
        "title": "Diagrama de Classes — Hierarquia de Usuários",
        "time": "1min15s",
        "cue": "AbstractUser no centro topo, 4 subclasses embaixo.",
        "speech": (
            "Esse é o diagrama de classes da nossa camada de "
            "domínio. No topo está a AbstractUser, uma classe "
            "abstrata que centraliza tudo o que é comum: id, "
            "nome, email, hash da senha, e métodos como block, "
            "approve e updateProfile. Embaixo, as quatro "
            "subclasses concretas — Admin, Donor, Volunteer e "
            "Receiver — herdam essa base. Repare nos elementos "
            "auxiliares: o UserFactory à direita, que olha o "
            "campo role do banco e devolve a subclasse certa; e "
            "o tipo Action à esquerda, que enumera as ações "
            "possíveis do sistema."
        ),
    },
    {
        "presenter": "Jhonatan",
        "title": "Por que Modelar Usuários Assim?",
        "time": "1min15s",
        "cue": "Três colunas: leia o título de cada uma.",
        "speech": (
            "Por que separar tudo isso? Pelos três pilares de "
            "orientação a objetos. Herança: a AbstractUser "
            "centraliza o comportamento comum, então a gente não "
            "duplica código. Polimorfismo: cada papel "
            "sobrescreve canPerform, então o middleware de "
            "autorização só chama user.canPerform de uma ação e "
            "o TypeScript despacha sozinho para a subclasse "
            "certa. E encapsulamento: o estado interno é "
            "protected, ninguém de fora consegue colocar "
            "isActive como falso direto — precisa passar pelo "
            "método block que valida invariantes."
        ),
    },
    {
        "presenter": "Jhonatan",
        "title": "Padrão Strategy",
        "time": "1min15s",
        "cue": "Diagrama à esquerda (interface + 3 estratégias), "
                "bullets à direita.",
        "speech": (
            "Primeiro padrão aplicado: Strategy. Quando uma "
            "doação é reservada, o doador precisa ser avisado. "
            "Mas o canal varia — pode ser email, in-app, push, "
            "ou só log em desenvolvimento. Em vez de espalhar "
            "ifs, cada canal é uma estratégia que implementa a "
            "mesma interface. O DonationsService recebe um "
            "NotificationService que delega para a estratégia "
            "atual. Pra adicionar SMS, basta criar uma nova "
            "classe — o service nem é tocado."
        ),
    },
    {
        "presenter": "Jhonatan",
        "title": "Padrão Observer",
        "time": "1min15s",
        "cue": "Subject e 3 observers à esquerda, bullets à direita.",
        "speech": (
            "Segundo padrão: Observer. Quando uma doação muda "
            "de status, vários efeitos colaterais podem "
            "acontecer — notificar o doador, gravar log de "
            "auditoria, invalidar o cache do dashboard. Em vez "
            "do service conhecer cada um deles, ele publica um "
            "evento no DonationEventBus. Os três observers que "
            "vêem aqui são inscritos lá no construtor e reagem "
            "de forma independente. Se um falhar, os outros "
            "continuam funcionando — isso a gente garante com "
            "try-catch dentro do notify."
        ),
    },
    {
        "presenter": "Jhonatan",
        "title": "Padrão Decorator",
        "time": "1min15s",
        "cue": "Component/ConcreteComponent + 3 decoradores; bullets "
                "à direita.",
        "speech": (
            "E o terceiro: Decorator. Aqui o desafio é envolver "
            "o repositório do Prisma com responsabilidades "
            "transversais. O AuditDecorator registra cada "
            "mutação no AuditLog. O TimingDecorator mede quantos "
            "milissegundos cada leitura levou. E o "
            "CacheDecorator memoriza os findOne, invalidando "
            "automaticamente quando há update. O legal é que "
            "todos compartilham a mesma interface, então a gente "
            "pode empilhar em qualquer ordem ou tirar uma "
            "camada sem mexer nas outras."
        ),
    },
    {
        "presenter": "Jhonatan",
        "title": "Cronograma — 7 quinzenas",
        "time": "1min",
        "cue": "Linha do tempo Q1 a Q7.",
        "speech": (
            "Esse é o nosso cronograma das sete quinzenas. Nas "
            "duas primeiras quinzenas, escolha do cenário e "
            "definição dos requisitos com o Jonas. Na terceira, "
            "modelagem do banco. Na quarta, o MVP e a "
            "apresentação parcial. Da quinta em diante, "
            "refinamento — validade de produtos, CPF de "
            "donatários, e a preparação do vídeo final. Estamos "
            "neste momento entregando a primeira versão dos "
            "artefatos, que corresponde a tudo o que vocês "
            "estão vendo aqui hoje."
        ),
    },
    {
        "presenter": "Jhonatan",
        "title": "Conclusão",
        "time": "45s",
        "cue": "Dois cards (esquerda: artefatos; direita: implementação).",
        "speech": (
            "Pra concluir: do lado da modelagem, entregamos os "
            "cinco diagramas obrigatórios, o diagrama de classes "
            "e os três diagramas dos padrões — cada um com PDF "
            "explicativo. Do lado da implementação, temos um "
            "backend NestJS funcional sobre PostgreSQL, com a "
            "hierarquia OO de usuários, Strategy aplicado nas "
            "notificações, Observer disparando nas mudanças de "
            "status, e Decorator pronto para envolver o "
            "repositório. O sistema substitui o caderno por um "
            "fluxo digital, rastreável e auditado."
        ),
    },
    {
        "presenter": "Jhonatan",
        "title": "Próximos Passos",
        "time": "45s",
        "cue": "Lista de 5 próximos passos.",
        "speech": (
            "E pra fechar, os próximos passos das quinzenas 5 a "
            "7: adicionar o campo de validade nas doações para "
            "ter alerta no dashboard, estender o cadastro de "
            "donatário com CPF e endereço, reimplementar o "
            "módulo de upload, substituir o canal Console por "
            "um provedor de email real, e ciclar feedback com "
            "o Jonas após a apresentação parcial."
        ),
    },
    {
        "presenter": "Grupo",
        "title": "Obrigado / Perguntas",
        "time": "30s",
        "cue": "Os três apresentadores se colocam visivelmente; "
                "convide perguntas.",
        "speech": (
            "É isso! Obrigado pela atenção. Estamos abertos a "
            "perguntas — qualquer um dos três pode responder, "
            "ou o Edilson, que nos ajudou no contato com a ONG."
        ),
    },
]


# =============================================================
# Aplica speaker notes ao .pptx
# =============================================================
def update_pptx():
    prs = Presentation(PPTX_PATH)
    if len(prs.slides) != len(SCRIPT):
        print(f"AVISO: pptx tem {len(prs.slides)} slides, "
              f"script tem {len(SCRIPT)} entradas.")
    for slide, entry in zip(prs.slides, SCRIPT):
        notes = slide.notes_slide.notes_text_frame
        notes.clear()

        speaker_name = PRESENTERS[entry["presenter"]]["name"]

        # 1ª linha: cabeçalho
        p0 = notes.paragraphs[0]
        r = p0.add_run()
        r.text = f"[{speaker_name}] · ~{entry['time']}"
        r.font.size = Pt(11)
        r.font.bold = True

        # Dica visual
        p1 = notes.add_paragraph()
        r = p1.add_run()
        r.text = f"Dica visual: {entry['cue']}"
        r.font.size = Pt(10)
        r.font.italic = True

        # Linha em branco + Fala
        notes.add_paragraph()
        p3 = notes.add_paragraph()
        r = p3.add_run()
        r.text = entry["speech"]
        r.font.size = Pt(12)

    prs.save(PPTX_PATH)
    print(f"-> notes atualizadas em {PPTX_PATH}")


# =============================================================
# Gera o PDF imprimível agrupado por apresentador
# =============================================================
def build_pdf():
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("h1", parent=styles["Heading1"],
                         textColor=NAVY, fontSize=22, spaceAfter=10,
                         alignment=TA_CENTER)
    sub = ParagraphStyle("sub", parent=styles["BodyText"],
                          textColor=GREY, fontSize=12, italic=True,
                          spaceAfter=14, alignment=TA_CENTER)
    h2 = ParagraphStyle("h2", parent=styles["Heading2"],
                         textColor=NAVY, fontSize=16, spaceAfter=8)
    h3 = ParagraphStyle("h3", parent=styles["Heading3"],
                         fontSize=12, spaceAfter=4)
    body = ParagraphStyle("body", parent=styles["BodyText"],
                            alignment=TA_JUSTIFY, fontSize=11,
                            spaceAfter=6, leading=15)
    speech = ParagraphStyle("speech", parent=styles["BodyText"],
                              alignment=TA_JUSTIFY, fontSize=12,
                              spaceAfter=8, leading=17,
                              leftIndent=10, rightIndent=10,
                              textColor=TEXT)
    cue = ParagraphStyle("cue", parent=styles["Italic"],
                          fontSize=10, spaceAfter=6,
                          textColor=GREY)
    meta = ParagraphStyle("meta", parent=styles["BodyText"],
                            fontSize=10, textColor=GREY, spaceAfter=6)

    doc = SimpleDocTemplate(PDF_PATH, pagesize=A4,
                             leftMargin=2 * cm, rightMargin=2 * cm,
                             topMargin=1.8 * cm, bottomMargin=1.5 * cm)
    story = []

    # ============= Capa =============
    story.append(Paragraph("Script de Apresentação", h1))
    story.append(Paragraph("Rede Solidária — Plataforma de Gestão de Doações",
                            sub))
    story.append(Spacer(1, 0.4 * cm))

    # Tabela resumo
    by_speaker = {"Ana": [], "Vitor": [], "Jhonatan": [], "Grupo": []}
    for i, e in enumerate(SCRIPT, 1):
        by_speaker[e["presenter"]].append((i, e))

    resumo_rows = [["Apresentador", "Slides", "Tempo aprox."]]
    for key in ["Ana", "Vitor", "Jhonatan", "Grupo"]:
        entries = by_speaker[key]
        if not entries:
            continue
        slide_nums = ", ".join(str(i) for i, _ in entries)
        # somar tempos
        total_s = 0
        for _, ent in entries:
            t = ent["time"]
            mins = secs = 0
            if "min" in t:
                parts = t.replace("s", "").split("min")
                mins = int(parts[0]) if parts[0].strip() else 0
                secs = int(parts[1]) if len(parts) > 1 and parts[1].strip() else 0
            else:
                secs = int(t.replace("s", "").strip())
            total_s += mins * 60 + secs
        m, s = total_s // 60, total_s % 60
        resumo_rows.append([
            PRESENTERS[key]["name"], slide_nums,
            f"~{m}min{s:02d}s" if s else f"~{m}min",
        ])
    # total geral
    total_grand = sum(
        (
            (int(e["time"].split("min")[0]) if "min" in e["time"] else 0) * 60
            + (
                int(e["time"].split("min")[1].replace("s", "").strip())
                if "min" in e["time"] and e["time"].split("min")[1].strip()
                else (int(e["time"].replace("s", "").strip())
                      if "min" not in e["time"] else 0)
            )
        )
        for e in SCRIPT
    )
    tm, ts = total_grand // 60, total_grand % 60
    resumo_rows.append(["Total", str(len(SCRIPT)),
                         f"~{tm}min{ts:02d}s"])

    t = Table(resumo_rows, colWidths=[5.5 * cm, 7.5 * cm, 3.5 * cm])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), NAVY),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("BACKGROUND", (0, -1), (-1, -1), GREY_L),
        ("FONTNAME", (0, -1), (-1, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("ALIGN", (1, 0), (-1, -1), "CENTER"),
        ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#cccccc")),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.append(t)
    story.append(Spacer(1, 0.5 * cm))

    story.append(Paragraph(
        "<b>Como usar este script:</b> ele está organizado por "
        "apresentador. Cada bloco começa com uma capa colorida com o "
        "nome do apresentador. Para cada slide há a <i>dica visual</i> "
        "(o que está aparecendo na tela e onde apontar) e a "
        "<i>fala completa</i> (texto pronto para ser lido em ritmo "
        "natural — não precisa decorar). O .pptx também recebeu as "
        "mesmas notas, então no Modo Apresentador do PowerPoint elas "
        "aparecem na sua tela enquanto a plateia vê só o slide.",
        body))

    story.append(PageBreak())

    # ============= Por apresentador =============
    for key in ["Ana", "Vitor", "Jhonatan", "Grupo"]:
        entries = by_speaker[key]
        if not entries:
            continue
        # Capa colorida do bloco
        p = PRESENTERS[key]
        cover_t = Table(
            [[Paragraph(
                f"<font color='white' size='28'><b>{p['name']}</b></font><br/>"
                f"<font color='white' size='14'><i>"
                f"{len(entries)} slides — clique para ver suas falas"
                f"</i></font>",
                styles["BodyText"])]],
            colWidths=[16 * cm], rowHeights=[7 * cm],
        )
        cover_t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), p["color"]),
            ("ALIGN", (0, 0), (-1, -1), "CENTER"),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (-1, -1), 20),
        ]))
        story.append(Spacer(1, 4 * cm))
        story.append(cover_t)
        story.append(PageBreak())

        # Slides desse apresentador
        for idx, e in entries:
            # cabeçalho do slide
            head = (
                f"<font color='{p['color'].hexval()}'><b>"
                f"Slide {idx} — {e['title']}</b></font>"
            )
            story.append(Paragraph(head, h2))
            story.append(Paragraph(
                f"<b>Tempo estimado:</b> {e['time']} "
                f"&nbsp;&nbsp;|&nbsp;&nbsp; <b>Apresentador:</b> {p['name']}",
                meta))
            # caixa pra dica visual
            cue_t = Table(
                [[Paragraph(
                    f"<b>Dica visual:</b> {e['cue']}",
                    cue)]],
                colWidths=[16 * cm],
            )
            cue_t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1), GREY_L),
                ("BOX", (0, 0), (-1, -1), 0.3, p["color"]),
                ("LEFTPADDING", (0, 0), (-1, -1), 8),
                ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]))
            story.append(cue_t)
            story.append(Spacer(1, 0.2 * cm))
            # fala
            story.append(Paragraph("<b>Fala:</b>", h3))
            story.append(Paragraph(f"“{e['speech']}”", speech))
            story.append(Spacer(1, 0.4 * cm))

        story.append(PageBreak())

    # ============= Dicas finais =============
    story.append(Paragraph("Dicas gerais para o dia da apresentação",
                            h2))
    tips = [
        "<b>Ensaiem em voz alta pelo menos uma vez juntos.</b> Não "
        "precisa decorar, mas o tempo só fica natural depois do "
        "primeiro ensaio.",
        "<b>Olhem para a plateia, não para a tela.</b> O .pptx tem as "
        "mesmas notas que estão neste PDF — use o Modo Apresentador.",
        "<b>Cuidado com a velocidade.</b> Quando a gente está "
        "nervoso, fala 30% mais rápido. Respire antes de cada slide.",
        "<b>Combine os gestos de transição.</b> Quem termina diz o "
        "nome do próximo (“agora passo pra o Vitor”), e o próximo "
        "começa cumprimentando (“obrigado, Ana”).",
        "<b>Se errar, siga em frente.</b> Se esquecer uma palavra, "
        "use os bullets do slide como apoio — eles foram desenhados "
        "pra isso.",
        "<b>Perguntas:</b> qualquer um dos três pode responder. Se "
        "for sobre código, é Jhonatan. Se for sobre a ONG, Ana "
        "ou Edilson. Se for sobre modelagem, Vitor.",
    ]
    for tip in tips:
        story.append(Paragraph("•  " + tip, body))

    doc.build(story)
    print(f"-> {PDF_PATH}")


if __name__ == "__main__":
    update_pptx()
    build_pdf()
