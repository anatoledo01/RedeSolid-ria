"""Gera a apresentação PowerPoint do projeto Rede Solidária."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


ROOT = "/Users/jhonatanrabelo/Documents/Fatec/Eng 3/Projeto/RedeSolid-ria"
DIAG = os.path.join(ROOT, "docs", "diagramas")
OUT = os.path.join(ROOT, "Apresentacao_RedeSolidaria.pptx")

# Paleta
NAVY      = RGBColor(0x1f, 0x3a, 0x5f)
NAVY_DARK = RGBColor(0x14, 0x29, 0x45)
GREEN     = RGBColor(0x10, 0xb9, 0x81)
GREEN_L   = RGBColor(0xd1, 0xfa, 0xe5)
AMBER     = RGBColor(0xb4, 0x53, 0x09)
AMBER_L   = RGBColor(0xfe, 0xf3, 0xc7)
GREY      = RGBColor(0x47, 0x55, 0x69)
GREY_L    = RGBColor(0xf1, 0xf5, 0xf9)
WHITE     = RGBColor(0xff, 0xff, 0xff)
TEXT      = RGBColor(0x0f, 0x17, 0x2a)

# Cores por apresentador (badge / faixas de transição)
PRESENTERS = {
    "Ana":      {"color": RGBColor(0xdb, 0x27, 0x77),
                 "name":  "Ana Toledo",
                 "role":  "Contexto, problema e solução"},
    "Vitor":    {"color": RGBColor(0x25, 0x63, 0xeb),
                 "name":  "Vitor Campos",
                 "role":  "Modelagem do sistema"},
    "Jhonatan": {"color": RGBColor(0x05, 0x96, 0x69),
                 "name":  "Jhonatan Rabelo",
                 "role":  "Implementação, padrões e conclusão"},
}


def add_blank_slide(prs):
    blank = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(blank)


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if line_color is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line_color
        rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    return rect


def add_text(slide, x, y, w, h, text, size=14, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT,
             font_name="Calibri", italic=False,
             v_anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = v_anchor

    lines = text.split("\n") if "\n" in text else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font_name
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=16, bullet="• "):
    """items: list of strings (or tuples (text, bold))."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    for i, it in enumerate(items):
        bold = False
        text = it
        if isinstance(it, tuple):
            text, bold = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = bullet + text
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = TEXT
        r.font.bold = bold
    return tb


def add_presenter_badge(slide, prs, presenter_key):
    """Chip canto superior direito com nome do apresentador."""
    p = PRESENTERS[presenter_key]
    w, h = Inches(2.6), Inches(0.45)
    x = prs.slide_width - w - Inches(0.35)
    y = Inches(0.20)
    add_rect(slide, x, y, w, h, p["color"])
    add_text(slide, x, y, w, h,
             f"●  {p['name']}",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)


def slide_header(slide, prs, title, kicker=None, presenter=None):
    """Barra superior padrão de cada slide de conteúdo."""
    # barra
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.85), NAVY)
    # kicker
    if kicker:
        add_text(slide, Inches(0.5), Inches(0.10), Inches(8), Inches(0.30),
                 kicker.upper(), size=11, color=GREEN_L, bold=True)
    # título
    add_text(slide, Inches(0.5), Inches(0.32), Inches(9.5), Inches(0.55),
             title, size=24, bold=True, color=WHITE)
    # badge do apresentador
    if presenter:
        add_presenter_badge(slide, prs, presenter)
    # rodapé
    add_rect(slide, Inches(0), prs.slide_height - Inches(0.35),
             prs.slide_width, Inches(0.35), GREY_L)
    add_text(slide, Inches(0.4), prs.slide_height - Inches(0.33),
             Inches(8), Inches(0.30),
             "Rede Solidária · Projeto de Curricularização · Eng. SW III",
             size=9, color=GREY, italic=True)


def transition_slide(prs, presenter_key, section_title, items):
    """Slide-divisor anunciando o próximo apresentador e o que vem."""
    p = PRESENTERS[presenter_key]
    s = add_blank_slide(prs)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, p["color"])
    # acento branco lateral
    add_rect(s, 0, 0, Inches(0.3), prs.slide_height, WHITE)
    # "Próximo apresentador"
    add_text(s, Inches(1.0), Inches(1.2), Inches(11), Inches(0.5),
             "PRÓXIMO APRESENTADOR",
             size=14, bold=True, color=WHITE, italic=False)
    # Nome grande
    add_text(s, Inches(1.0), Inches(1.8), Inches(11), Inches(1.4),
             p["name"], size=56, bold=True, color=WHITE)
    # Seção
    add_text(s, Inches(1.0), Inches(3.4), Inches(11), Inches(0.6),
             section_title, size=26, color=WHITE, italic=True)
    # Itens que ele vai cobrir
    add_text(s, Inches(1.0), Inches(4.4), Inches(11), Inches(0.4),
             "O QUE VEM A SEGUIR",
             size=11, bold=True, color=WHITE)
    tb = s.shapes.add_textbox(Inches(1.0), Inches(4.8),
                                Inches(11), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, t in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(6)
        para.alignment = PP_ALIGN.LEFT
        r = para.add_run()
        r.text = "→  " + t
        r.font.name = "Calibri"
        r.font.size = Pt(20)
        r.font.color.rgb = WHITE
        r.font.bold = False
    return s


def add_image_centered(slide, prs, image_path,
                        top=Inches(1.1), bottom=Inches(0.5),
                        side_margin=Inches(0.6)):
    """Adiciona imagem centralizada respeitando margens do slide."""
    max_w = prs.slide_width - side_margin * 2
    max_h = prs.slide_height - top - bottom
    pic = slide.shapes.add_picture(image_path, side_margin, top,
                                    width=max_w)
    # se a imagem ficou alta demais, recalcula pela altura
    if pic.height > max_h:
        slide.shapes._spTree.remove(pic._element)
        ratio = max_h / pic.height
        new_w = int(pic.width * ratio)
        new_h = max_h
        x = (prs.slide_width - new_w) // 2
        pic = slide.shapes.add_picture(image_path, x, top,
                                        width=new_w, height=new_h)
    else:
        # centra horizontalmente
        x = (prs.slide_width - pic.width) // 2
        pic.left = x
    return pic


# ====================================================================
# CONSTRUÇÃO DA APRESENTAÇÃO
# ====================================================================
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height

    # ---------- SLIDE 1 — CAPA ----------
    s = add_blank_slide(prs)
    # Background dividido
    add_rect(s, 0, 0, W, H, NAVY)
    # Acento verde
    add_rect(s, 0, Inches(6.5), W, Inches(1), GREEN)
    # Marca/heart símbolo simulado
    add_rect(s, Inches(1.0), Inches(1.0), Inches(0.55), Inches(0.55), GREEN)
    add_text(s, Inches(1.65), Inches(0.95), Inches(8), Inches(0.5),
             "Rede Solidária", size=18, bold=True, color=WHITE)
    add_text(s, Inches(1.65), Inches(1.32), Inches(8), Inches(0.3),
             "Conectando solidariedade com quem mais precisa",
             size=10, color=GREEN_L, italic=True)
    # Título central
    add_text(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.3),
             "Plataforma de Gestão\nde Doações",
             size=54, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Subtítulo
    add_text(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.5),
             "Projeto de Curricularização — Engenharia de Software III",
             size=20, color=GREEN_L)
    add_text(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.4),
             "Parceria com a ONG Cidade Social de Mogi",
             size=16, color=WHITE, italic=True)
    # Equipe
    add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.4),
             "EQUIPE",
             size=11, bold=True, color=GREEN_L)
    add_text(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
             "Jhonatan Rabelo · Vitor Campos · Ana Toledo · Edilson Junior",
             size=16, color=WHITE)
    # data no acento verde
    add_text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
             "Mogi das Cruzes · 2026",
             size=14, bold=True, color=NAVY)

    # ---------- SLIDE 2 — AGENDA + DIVISÃO ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Agenda e Divisão", "Roteiro", presenter="Ana")
    add_text(s, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.5),
             "Roteiro da apresentação e como dividimos a fala entre nós",
             size=14, color=GREY, italic=True)
    blocks = [
        ("Ana", "Bloco 1  —  História do projeto",
         ["O problema na Cidade Social de Mogi",
          "A solução proposta: Rede Solidária"]),
        ("Vitor", "Bloco 2  —  Modelagem do sistema",
         ["Stack técnica",
          "Diagramas IDEF0 N0/N1, Casos de Uso, DFD e DER"]),
        ("Jhonatan", "Bloco 3  —  Implementação e fechamento",
         ["Modelagem OO + Diagrama de Classes",
          "Padrões aplicados: Strategy, Observer e Decorator",
          "Cronograma, conclusão e próximos passos"]),
    ]
    top = Inches(1.7)
    for pk, title, sub_items in blocks:
        p = PRESENTERS[pk]
        # faixa colorida com nome
        add_rect(s, Inches(0.8), top, Inches(2.8), Inches(1.55),
                 p["color"])
        add_text(s, Inches(0.8), top + Inches(0.25),
                 Inches(2.8), Inches(0.35),
                 p["name"], size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.8), top + Inches(0.65),
                 Inches(2.8), Inches(0.3),
                 "●●●", size=10, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.8), top + Inches(0.95),
                 Inches(2.8), Inches(0.5),
                 p["role"], size=10, color=WHITE,
                 italic=True, align=PP_ALIGN.CENTER)
        # corpo
        add_rect(s, Inches(3.7), top, Inches(8.9), Inches(1.55),
                 GREY_L, p["color"])
        add_text(s, Inches(3.9), top + Inches(0.15),
                 Inches(8.5), Inches(0.4),
                 title, size=15, bold=True, color=p["color"])
        bullets_tb = s.shapes.add_textbox(Inches(3.9), top + Inches(0.55),
                                           Inches(8.5), Inches(0.95))
        bf = bullets_tb.text_frame
        bf.word_wrap = True
        for i, it in enumerate(sub_items):
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            r = para.add_run()
            r.text = "•  " + it
            r.font.size = Pt(12)
            r.font.color.rgb = TEXT
        top += Inches(1.65)

    # ---------- SLIDE 3 — O PROBLEMA (contexto) ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "O Problema", "Parte 1 — Contexto",
                 presenter="Ana")
    add_text(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
             "Cidade Social de Mogi — Av. São Paulo, 33, Jardim Armênia",
             size=18, bold=True, color=NAVY)
    add_text(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.5),
             "ONG que monta e distribui cestas básicas para famílias em "
             "vulnerabilidade",
             size=14, color=GREY, italic=True)
    # Box destaque
    add_rect(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.7),
             GREEN_L, GREEN)
    add_text(s, Inches(1.0), Inches(2.75), Inches(11.3), Inches(0.45),
             "Entrevista com Jonas Oliveira (stakeholder)",
             size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(3.25), Inches(11.3), Inches(1.1),
             "“Hoje a gente perde produto por validade vencida, não consegue "
             "saber direito quem doou o quê e nem que cesta foi para qual "
             "família. Tudo no caderno.”",
             size=14, italic=True, color=TEXT)
    # Estatística simples
    add_text(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.5),
             "O processo atual é manual, em papel — escala mal e perde "
             "rastreabilidade.",
             size=14, color=TEXT)

    # ---------- SLIDE 4 — O PROBLEMA (dores específicas) ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "O Problema", "Parte 2 — Dores específicas",
                 presenter="Ana")
    add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.5),
             "Necessidades levantadas com a ONG", size=18, bold=True,
             color=NAVY)

    # 3 colunas tipo card
    col_w = Inches(3.8)
    col_h = Inches(4.3)
    top = Inches(1.9)
    gap = Inches(0.25)
    cards = [
        ("Controle de Estoque",
         "Monitorar validade dos produtos para montar cestas "
         "nutricionalmente equilibradas com o que está em dia."),
        ("Cadastro de Pessoas",
         "Manter registro completo de doadores e donatários "
         "(CPF, nome, contato)."),
        ("Rastreamento de Cestas",
         "Saber quando cada cesta foi montada, sua validade e para "
         "qual família foi entregue."),
    ]
    x = Inches(0.8)
    for title, desc in cards:
        add_rect(s, x, top, col_w, col_h, WHITE, NAVY)
        # faixa colorida no topo do card
        add_rect(s, x, top, col_w, Inches(0.5), NAVY)
        add_text(s, x + Inches(0.2), top + Inches(0.07),
                 col_w - Inches(0.4), Inches(0.4),
                 title, size=14, bold=True, color=WHITE)
        add_text(s, x + Inches(0.25), top + Inches(0.75),
                 col_w - Inches(0.5), col_h - Inches(0.9),
                 desc, size=14, color=TEXT)
        x += col_w + gap

    add_text(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
             "Resumo: a ONG precisa transformar planilhas e cadernos "
             "em um processo digital, simples e auditável.",
             size=14, color=GREY, italic=True)

    # ---------- SLIDE 5 — A SOLUÇÃO (visão geral) ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "A Solução", "Visão geral", presenter="Ana")
    add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.6),
             "Rede Solidária — plataforma web full-stack",
             size=22, bold=True, color=NAVY)
    add_text(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.5),
             "Conecta quatro papéis em torno do ciclo da doação",
             size=14, color=GREY, italic=True)

    # 4 papéis em cards
    roles = [
        ("Doador",   "Publica itens, anexa fotos,\nacompanha o destino"),
        ("Recebedor", "Busca o que precisa\ne reserva a doação"),
        ("Voluntário", "Aceita entregas\ne atualiza status"),
        ("Administrador", "Aprova usuários, modera,\nvê dashboard e logs"),
    ]
    col_w = Inches(2.9)
    col_h = Inches(2.6)
    top = Inches(2.4)
    x = Inches(0.8)
    gap = Inches(0.18)
    for nm, dsc in roles:
        add_rect(s, x, top, col_w, col_h, GREEN_L, GREEN)
        add_rect(s, x, top, col_w, Inches(0.55), GREEN)
        add_text(s, x, top + Inches(0.1), col_w, Inches(0.45),
                 nm, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), top + Inches(0.75),
                 col_w - Inches(0.4), col_h - Inches(0.85),
                 dsc, size=13, color=TEXT, align=PP_ALIGN.CENTER)
        x += col_w + gap

    add_text(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5),
             "Fluxo: Doador publica  →  Recebedor reserva  →  "
             "Voluntário entrega  →  ambos avaliam",
             size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.4),
             "Tudo auditado, com notificações automáticas em cada "
             "mudança de status.",
             size=13, color=GREY, italic=True)

    # ---------- TRANSIÇÃO  →  VITOR ----------
    transition_slide(prs, "Vitor", "Bloco 2  —  Modelagem do Sistema",
                     ["Stack tecnológica que escolhemos",
                      "Os 5 diagramas que descrevem o sistema:",
                      "    IDEF0 N0 e N1, Casos de Uso, DFD e DER"])

    # ---------- SLIDE — STACK ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Stack Tecnológica", "Desenvolvimento",
                 presenter="Vitor")
    layers = [
        ("Frontend",
         "Next.js 15 · React 19 · TypeScript · TailwindCSS v4 · "
         "Zustand · React Query · Framer Motion"),
        ("Backend",
         "NestJS 11 · TypeScript · Prisma ORM 6 · JWT + Passport · "
         "Helmet · Throttler · Swagger/OpenAPI"),
        ("Banco",
         "PostgreSQL 16  (Docker Compose, porta 5433)"),
        ("Dev tooling",
         "Docker · npm · ts-node · ESLint"),
    ]
    top = Inches(1.4)
    for nm, dsc in layers:
        add_rect(s, Inches(0.8), top, Inches(11.7), Inches(1.05), GREY_L)
        add_rect(s, Inches(0.8), top, Inches(2.4), Inches(1.05), NAVY)
        add_text(s, Inches(0.8), top + Inches(0.3),
                 Inches(2.4), Inches(0.45),
                 nm, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(3.4), top + Inches(0.3),
                 Inches(9), Inches(0.6),
                 dsc, size=14, color=TEXT,
                 v_anchor=MSO_ANCHOR.MIDDLE)
        top += Inches(1.2)

    add_text(s, Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.3),
             "Backend: http://localhost:3001/api/docs  •  "
             "Frontend: http://localhost:3000",
             size=11, color=GREEN, italic=True, bold=True)

    # ---------- SLIDE 7 — IDEF0 N0 ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Visão de Sistema  —  IDEF0 N0",
                 "Modelagem · Diagramas", presenter="Vitor")
    add_image_centered(s, prs, os.path.join(DIAG, "IDEF0_N0.png"),
                        top=Inches(1.0), bottom=Inches(0.45))

    # ---------- SLIDE 8 — IDEF0 N1 ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Decomposição em 5 Atividades  —  IDEF0 N1",
                 "Modelagem · Diagramas", presenter="Vitor")
    add_image_centered(s, prs, os.path.join(DIAG, "IDEF0_N1.png"),
                        top=Inches(1.0), bottom=Inches(0.45))

    # ---------- SLIDE 9 — Casos de Uso ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Casos de Uso",
                 "Modelagem · Diagramas", presenter="Vitor")
    add_image_centered(s, prs, os.path.join(DIAG, "CasosDeUso.png"),
                        top=Inches(1.0), bottom=Inches(0.45))

    # ---------- SLIDE 10 — DFD ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Fluxo de Dados (DFD Nível 1)",
                 "Modelagem · Diagramas", presenter="Vitor")
    add_image_centered(s, prs, os.path.join(DIAG, "DFD.png"),
                        top=Inches(1.0), bottom=Inches(0.45))

    # ---------- SLIDE 11 — DER ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Modelo de Dados (DER)",
                 "Modelagem · Diagramas", presenter="Vitor")
    add_image_centered(s, prs, os.path.join(DIAG, "DER.png"),
                        top=Inches(1.0), bottom=Inches(0.45))

    # ---------- TRANSIÇÃO  →  JHONATAN ----------
    transition_slide(prs, "Jhonatan",
                     "Bloco 3  —  Implementação e fechamento",
                     ["Modelagem OO da hierarquia de usuários",
                      "Padrões aplicados: Strategy, Observer e Decorator",
                      "Cronograma do projeto, conclusão e próximos passos"])

    # ---------- SLIDE — Modelagem OO ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Modelagem OO  —  Hierarquia de Usuários",
                 "Implementação · OO", presenter="Jhonatan")
    add_image_centered(s, prs, os.path.join(DIAG, "DiagramaClasses.png"),
                        top=Inches(1.0), bottom=Inches(0.45))

    # ---------- SLIDE 13 — OO conceitos (texto) ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Por que Modelar Usuários Assim?",
                 "Implementação · OO", presenter="Jhonatan")
    # 3 colunas: herança / polimorfismo / encapsulamento
    cards = [
        ("HERANÇA",
         "AbstractUser centraliza os atributos comuns e o comportamento "
         "compartilhado (block, approve, updateProfile).\n\n"
         "Admin, Donor, Volunteer e Receiver herdam essa base.",
         NAVY, NAVY),
        ("POLIMORFISMO",
         "Cada papel sobrescreve canPerform(action) e getDashboardRoute().\n\n"
         "Quem chama user.canPerform('delivery:accept') não sabe — nem "
         "precisa saber — qual subclasse está respondendo.",
         GREEN, GREEN_L),
        ("ENCAPSULAMENTO",
         "Estado privado/protegido — mudanças passam por métodos que "
         "validam invariantes (ex.: não bloquear já bloqueado).\n\n"
         "Senha nunca aparece em toPublicJSON().",
         AMBER, AMBER_L),
    ]
    col_w = Inches(3.95)
    col_h = Inches(5.0)
    top = Inches(1.4)
    x = Inches(0.6)
    gap = Inches(0.15)
    for title, desc, edge, face in cards:
        add_rect(s, x, top, col_w, col_h, face, edge)
        add_rect(s, x, top, col_w, Inches(0.55), edge)
        add_text(s, x, top + Inches(0.1), col_w, Inches(0.4),
                 title, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), top + Inches(0.8),
                 col_w - Inches(0.5), col_h - Inches(0.9),
                 desc, size=13, color=TEXT)
        x += col_w + gap

    add_text(s, Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.4),
             "Arquivo: backend/src/domain/users/user.entity.ts",
             size=11, color=GREY, italic=True, align=PP_ALIGN.CENTER)

    # ---------- SLIDE 14 — Strategy ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Padrão Strategy  —  Notificações",
                 "Padrões de Projeto", presenter="Jhonatan")
    # imagem à esquerda
    pic_w = Inches(7.6)
    pic = s.shapes.add_picture(os.path.join(DIAG, "Padrao_Strategy.png"),
                                Inches(0.4), Inches(1.05),
                                width=pic_w)
    # ajusta altura máxima
    max_h = Inches(5.7)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = max_h

    # texto à direita
    tx = Inches(8.2)
    tw = Inches(4.9)
    add_text(s, tx, Inches(1.1), tw, Inches(0.5),
             "Trocar canal sem mudar o cliente",
             size=18, bold=True, color=NAVY)
    add_bullets(s, tx, Inches(1.7), tw, Inches(4.5),
                [
                    "Cada canal é uma estratégia: Console (dev), "
                    "Email (prod), In-app.",
                    "NotificationService recebe a estratégia no construtor.",
                    "DonationsService.reserve() apenas chama notify(...).",
                    "Adicionar SMS ou Push = nova classe; nada muda no "
                    "service.",
                ], size=13)
    add_text(s, tx, Inches(6.5), tw, Inches(0.4),
             "Arquivo: domain/notifications/notification.strategy.ts",
             size=10, color=GREY, italic=True)

    # ---------- SLIDE 15 — Observer ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Padrão Observer  —  Eventos de Doação",
                 "Padrões de Projeto", presenter="Jhonatan")
    pic = s.shapes.add_picture(os.path.join(DIAG, "Padrao_Observer.png"),
                                Inches(0.4), Inches(1.05),
                                width=Inches(7.6))
    max_h = Inches(5.7)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = max_h
    tx = Inches(8.2)
    tw = Inches(4.9)
    add_text(s, tx, Inches(1.1), tw, Inches(0.5),
             "Um evento, vários efeitos colaterais",
             size=18, bold=True, color=NAVY)
    add_bullets(s, tx, Inches(1.7), tw, Inches(4.5),
                [
                    "DonationEventBus emite eventos de mudança de status.",
                    "Observers inscritos reagem independentemente:",
                    "    – NotifyDonorObserver  (email)",
                    "    – AuditTrailObserver   (AuditLog)",
                    "    – DashboardCacheObserver  (invalida cache)",
                    "Falha de um não derruba os outros.",
                ], size=13)
    add_text(s, tx, Inches(6.5), tw, Inches(0.4),
             "Arquivo: domain/events/donation-event-bus.ts",
             size=10, color=GREY, italic=True)

    # ---------- SLIDE 16 — Decorator ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Padrão Decorator  —  Repositório de Doações",
                 "Padrões de Projeto", presenter="Jhonatan")
    pic = s.shapes.add_picture(os.path.join(DIAG, "Padrao_Decorator.png"),
                                Inches(0.4), Inches(1.05),
                                width=Inches(7.6))
    max_h = Inches(5.7)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = max_h
    tx = Inches(8.2)
    tw = Inches(4.9)
    add_text(s, tx, Inches(1.1), tw, Inches(0.5),
             "Empilhar responsabilidades transversais",
             size=18, bold=True, color=NAVY)
    add_bullets(s, tx, Inches(1.7), tw, Inches(4.5),
                [
                    "PrismaDonationRepository fala com o banco.",
                    "Decoradores envolvem com novas capacidades:",
                    "    – AuditDecorator   (registra mutações)",
                    "    – TimingDecorator  (mede ms de leitura)",
                    "    – CacheDecorator   (memoiza findOne)",
                    "Compor / remover camada sem tocar nas outras.",
                ], size=13)
    add_text(s, tx, Inches(6.5), tw, Inches(0.4),
             "Arquivo: domain/decorators/donation-repository.ts",
             size=10, color=GREY, italic=True)

    # ---------- SLIDE 17 — Cronograma ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Cronograma  —  7 quinzenas",
                 "Desenvolvimento", presenter="Jhonatan")
    add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.5),
             "Plano de ação aprovado e cumprido entre fev/2026 e jun/2026",
             size=14, color=GREY, italic=True)
    rows = [
        ("Q1", "23/02 – 08/03",
         "Levantamento de ideias e escolha do cenário (ONG sugerida pelo "
         "Edilson)."),
        ("Q2", "09/03 – 23/03",
         "Reunião com Jonas, definição dos requisitos e plano de ação."),
        ("Q3", "28/03 – 07/04",
         "Definição do título e modelagem do banco PostgreSQL."),
        ("Q4", "08/04 – 22/04",
         "MVP (CRUD) + relatório parcial + apresentação inicial à ONG."),
        ("Q5", "23/04 – 07/05",
         "Refatoração, alertas de validade, revisão de cadastros (CPF)."),
        ("Q6", "08/05 – 22/05",
         "Testes com usuários reais, roteiro e gravação do vídeo."),
        ("Q7", "23/05 – 06/06",
         "Edição do vídeo, formatação do relatório final, entrega."),
    ]
    top = Inches(1.85)
    row_h = Inches(0.7)
    for q, dt, desc in rows:
        # faixa Q
        add_rect(s, Inches(0.8), top, Inches(0.7), row_h, NAVY)
        add_text(s, Inches(0.8), top, Inches(0.7), row_h,
                 q, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
        # data
        add_rect(s, Inches(1.5), top, Inches(2.2), row_h, GREY_L)
        add_text(s, Inches(1.5), top, Inches(2.2), row_h,
                 dt, size=11, color=NAVY, align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE, bold=True)
        # descrição
        add_rect(s, Inches(3.7), top, Inches(8.8), row_h, WHITE, GREY)
        add_text(s, Inches(3.85), top, Inches(8.6), row_h,
                 desc, size=12, color=TEXT,
                 v_anchor=MSO_ANCHOR.MIDDLE)
        top += row_h

    # ---------- SLIDE 18 — CONCLUSÃO ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Conclusão", "Encerramento", presenter="Jhonatan")
    add_text(s, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.6),
             "O que entregamos",
             size=22, bold=True, color=NAVY)

    # 2 colunas
    col_w = Inches(5.9)
    col_h = Inches(4.5)
    top = Inches(1.9)
    # esquerda: artefatos
    add_rect(s, Inches(0.8), top, col_w, col_h, GREEN_L, GREEN)
    add_rect(s, Inches(0.8), top, col_w, Inches(0.55), GREEN)
    add_text(s, Inches(0.8), top + Inches(0.1), col_w, Inches(0.4),
             "Artefatos de modelagem", size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(1.0), top + Inches(0.75),
                col_w - Inches(0.4), col_h - Inches(0.9),
                [
                    "5 diagramas: IDEF0 N0/N1, Casos de Uso, DFD, DER",
                    "1 diagrama de classes UML",
                    "3 diagramas de padrões de projeto",
                    "Cada um em PDF + PNG, com PDF explicativo",
                ], size=14)

    # direita: implementação
    add_rect(s, Inches(6.9), top, col_w, col_h, AMBER_L, AMBER)
    add_rect(s, Inches(6.9), top, col_w, Inches(0.55), AMBER)
    add_text(s, Inches(6.9), top + Inches(0.1), col_w, Inches(0.4),
             "Implementação real", size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(7.1), top + Inches(0.75),
                col_w - Inches(0.4), col_h - Inches(0.9),
                [
                    "Backend NestJS funcional sobre PostgreSQL",
                    "Hierarquia OO no domínio (User + 4 papéis)",
                    "Strategy aplicado nas notificações",
                    "Observer disparando em mudanças de status",
                    "Decorator pronto para envolver o repositório",
                ], size=14)

    add_text(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.5),
             "O sistema substitui o caderno por um fluxo digital, "
             "rastreável e auditado.",
             size=14, color=NAVY, bold=True, italic=True,
             align=PP_ALIGN.CENTER)

    # ---------- SLIDE 19 — PRÓXIMOS PASSOS ----------
    s = add_blank_slide(prs)
    slide_header(s, prs, "Próximos Passos", "Encerramento",
                 presenter="Jhonatan")
    add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.6),
             "Quinzenas 5 a 7 — caminho até a versão final",
             size=20, bold=True, color=NAVY)
    nexts = [
        ("Validade dos produtos",
         "Adicionar campo expiresAt nas doações e alertar no dashboard "
         "quando estiver próximo do vencimento."),
        ("CPF dos donatários",
         "Estender o cadastro de Usuário com CPF e endereço obrigatórios "
         "para o papel Recebedor."),
        ("Upload de imagens",
         "Reimplementar o UploadsModule para o doador anexar fotos reais "
         "do item."),
        ("Notificação por e-mail real",
         "Substituir ConsoleNotificationStrategy por EmailNotificationStrategy "
         "(SendGrid ou Amazon SES)."),
        ("Validação com a ONG",
         "Ciclo de feedback com Jonas Oliveira após a apresentação parcial."),
    ]
    top = Inches(1.9)
    row_h = Inches(0.85)
    for nm, dsc in nexts:
        add_rect(s, Inches(0.8), top, Inches(3.3), row_h, NAVY)
        add_text(s, Inches(0.9), top, Inches(3.1), row_h,
                 nm, size=13, bold=True, color=WHITE,
                 v_anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(4.1), top, Inches(8.4), row_h, GREY_L, GREY)
        add_text(s, Inches(4.25), top, Inches(8.2), row_h,
                 dsc, size=12, color=TEXT,
                 v_anchor=MSO_ANCHOR.MIDDLE)
        top += row_h + Inches(0.1)

    # ---------- SLIDE 20 — OBRIGADO ----------
    s = add_blank_slide(prs)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, Inches(6.5), W, Inches(1), GREEN)
    add_text(s, Inches(0), Inches(2.7), W, Inches(1.2),
             "Obrigado!",
             size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.0), W, Inches(0.6),
             "Perguntas?",
             size=24, color=GREEN_L, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(0), Inches(5.1), W, Inches(0.4),
             "Ana Toledo  ·  Vitor Campos  ·  Jhonatan Rabelo",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(5.55), W, Inches(0.4),
             "(com Edilson Junior — contato com a ONG)",
             size=12, color=GREEN_L, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(0), Inches(6.0), W, Inches(0.4),
             "FATEC  —  Engenharia de Software III",
             size=14, color=GREEN_L, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(0), Inches(6.75), W, Inches(0.5),
             "github.com/redesolidaria  ·  parceria com a Cidade Social de Mogi",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"-> {OUT}")
    print(f"   {len(prs.slides)} slides gerados")


if __name__ == "__main__":
    build()
